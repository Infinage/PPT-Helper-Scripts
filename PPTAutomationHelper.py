"""
Helper functions:

    1. Modify PPT Links
    2. Toggle update links popup
    3. Refresh excel plots (deps: xlwings)
    4. Refresh PPT links (deps: pywin32)
    5. Extract M Query file from excel
    6. Update M Query files inside excel
    7. Extract PPT to zip file
    8. Rezip extracted zip to PPT
    9. Update PPT Plot 'cache' (deps: lxml)
    10. Update excel embedded into PPT (deps: xlwings)
    11. Update PPT table (deps: python-pptx)

"""

import glob
import urllib
import zipfile
import os
import shutil
import tempfile
import datetime as dt
import win32com.client
import lxml.etree
import pptx
import pandas as pd
import xlwings as xw
import gzip
import base64
import re
import io

def _create_file_backup(path: str) -> str:
    """
    Given an input path, it adds a timestamp portion to the filename and stores a backup
    """
    existing_path = list(os.path.split(path))
    file_extension_pos = existing_path[-1].rfind(".")
    existing_path[-1] = f"{existing_path[-1][:file_extension_pos]}-{int(dt.datetime.today().timestamp())}{existing_path[-1][file_extension_pos:]}"
    backup_path = os.path.join(*existing_path)
    shutil.copy2(path, backup_path)
    return backup_path

def _search_and_replace(search_for: str, replace_with: str, file_paths: list[str]) -> int:
    """
    Search for the given string and replace with provided string on all given file paths.
    """

    print(f"'{search_for}' would be replaced with '{replace_with}'")
    counts = 0
    for file in file_paths:
        with open(file, "r", encoding="utf-8") as f:
            temp = f.read()

        counts += temp.count(search_for)
        temp = temp.replace(search_for, replace_with)

        with open(file, "w", encoding="utf-8") as f:
            f.write(temp)

    print(f"{counts} occurances were replaced.")
    return counts

def _extract_excel_datamashup(datamashup_byte: bytes, output_extract_path: str) -> int:
    """
    Given a datamashup from excel containing M queries, extracts the M query portion to extract path specified.
    Note that the datamashup is by itself an zip file.

    Contains logic specific to reading MS-QDEFF format zip

    https://community.fabric.microsoft.com/t5/Desktop/Change-pbix-data-source-programmatically/m-p/422128
    """

    byte_dict = {}
    with io.BytesIO(datamashup_byte) as dmf:
        # Parse bytes based on the MS-QDEFF documentation
        byte_dict['version'] = dmf.read(4)
        byte_dict['pkg_parts_len'] = dmf.read(4)
        byte_dict['pkg_parts'] = dmf.read(int.from_bytes(byte_dict['pkg_parts_len'], byteorder='little'))
        byte_dict['perm_len'] = dmf.read(4)
        byte_dict['perm_var'] = dmf.read(int.from_bytes(byte_dict['perm_len'], byteorder='little'))
        byte_dict['meta_len'] = dmf.read(4)
        byte_dict['meta_var'] = dmf.read(int.from_bytes(byte_dict['meta_len'], byteorder='little'))
        byte_dict['perm_bind_len'] = dmf.read(4)
        byte_dict['perm_bind_var'] = dmf.read(int.from_bytes(byte_dict['perm_bind_len'], byteorder='little'))

        # Keep track of number of files extracted
        count = 0
        with zipfile.ZipFile(io.BytesIO(byte_dict['pkg_parts'])) as zfd:
            for zfd_item in zfd.infolist():
                if re.match("Formulas/Section\d+\.m", zfd_item.filename):
                    zfd_item.filename = os.path.basename(zfd_item.filename)
                    zfd.extract(zfd_item, output_extract_path)
                    count += 1

        return count

def _update_excel_datamashup(datamashup_byte: bytes, mq_paths: dict[str, str]) -> bytes:
    """
    This function is used by the update M Query function.
    Contains logic specific to reading the MS-QDEFF format zip and updating existing data mashup.
    `mq_paths` provided contains the file paths to '.m' files to replace with.

    https://community.fabric.microsoft.com/t5/Desktop/Change-pbix-data-source-programmatically/m-p/422128
    """
    byte_dict = {}
    with io.BytesIO(datamashup_byte) as dmf:
        # Parse bytes based on the MS-QDEFF documentation
        byte_dict['version'] = dmf.read(4)
        byte_dict['pkg_parts_len'] = dmf.read(4)
        byte_dict['pkg_parts'] = dmf.read(int.from_bytes(byte_dict['pkg_parts_len'], byteorder='little'))
        byte_dict['perm_len'] = dmf.read(4)
        byte_dict['perm_var'] = dmf.read(int.from_bytes(byte_dict['perm_len'], byteorder='little'))
        byte_dict['meta_len'] = dmf.read(4)
        byte_dict['meta_var'] = dmf.read(int.from_bytes(byte_dict['meta_len'], byteorder='little'))
        byte_dict['perm_bind_len'] = dmf.read(4)
        byte_dict['perm_bind_var'] = dmf.read(int.from_bytes(byte_dict['perm_bind_len'], byteorder='little'))

        # Zip contents would be processed and stored into a temp buffer
        updated_zip_buffer = io.BytesIO()
        with zipfile.ZipFile(io.BytesIO(byte_dict['pkg_parts'])) as zfr, zipfile.ZipFile(updated_zip_buffer, mode='w', compression=zipfile.ZIP_DEFLATED) as zfw:
            for zfr_item in zfr.infolist():
                fname = os.path.basename(zfr_item.filename)
                if re.match('Formulas/Section\d+\.m', zfr_item.filename) and fname in mq_paths:
                    zfw.writestr(zfr_item, open(mq_paths[fname]).read())
                else:
                    zfw.writestr(zfr_item, zfr.open(zfr_item).read())

        # Update zip buffer contents to dict we have saved so far
        updated_zip_buffer.seek(0)
        byte_dict['pkg_parts'] = updated_zip_buffer.read()
        byte_dict['pkg_parts_len'] = len(byte_dict['pkg_parts']).to_bytes(4, byteorder='little')
        updated_zip_buffer.close()

        # Write out the revised datamashup file (in correct order, from ordered list)
        new_mashup_buffer = io.BytesIO()
        for b in ['version', 'pkg_parts_len', 'pkg_parts', 'perm_len', 'perm_var', 'meta_len', 'meta_var', 'perm_bind_len', 'perm_bind_var']:
            new_mashup_buffer.write(byte_dict[b])

        # We would be returning this updated mashup to overwrite the existing mashup content
        new_mashup_buffer.seek(0)
        new_mashup = new_mashup_buffer.read()
        new_mashup_buffer.close()
        return new_mashup

def modify_ppt_links(ppt_path: str, search_str: str, replace_with: str, overwrite: bool = False) -> None:
    """
    Helper function to modify PPT links to embedded excel objects.

    Please note that `replace_with` must be an absolute path only.
    If `replace_with` is not provided, the links are "broken".
    """

    # Create backup file before proceeding
    if not overwrite:
        backup_path = _create_file_backup(ppt_path)

    # Create a temp directory to store our intermediate stuff
    with tempfile.TemporaryDirectory() as tmpdir:
        zfw = zipfile.ZipFile(f"{tmpdir}/temp.pptx", "w", compression=zipfile.ZIP_DEFLATED)
        with zipfile.ZipFile(ppt_path) as zfr:
            for ITEM in zfr.infolist():
                if ITEM.filename.startswith("ppt/"):
                    zfr.extract(ITEM.filename, tmpdir)
                else:
                    zfw.writestr(ITEM, zfr.open(ITEM).read(), compress_type=zipfile.ZIP_DEFLATED)

        # Replace all spaces with '%20' in the search and replace strings
        REL_FILES = glob.glob(f"{tmpdir}/ppt/**/_rels/*.rels")
        search_str = "file:///" + search_str.replace(" ", "%20").replace("/", "\\")
        replace_with = "file:///" + replace_with.replace(" ", "%20").replace("/", "\\")
        _search_and_replace(search_str, replace_with, REL_FILES)

        # Add the PPT folder we have updated to the archive
        for FILE_NAME in glob.glob(f"{tmpdir}\\ppt\\**", recursive=True):
            if os.path.isfile(FILE_NAME):
                FNP = FILE_NAME.split("\\")
                ZFNP = os.path.join(*FNP[FNP.index("ppt"):])
                zfw.write(FILE_NAME, ZFNP, compress_type=zipfile.ZIP_DEFLATED)

        # Close to save it
        zfw.close()

        # Replace the existing PPT
        shutil.move(f"{tmpdir}/temp.pptx", ppt_path)

def toggle_update_links_popup(ppt_path: str, auto_update: bool = False, overwrite: bool = False) -> None:
    """
    Helper function to toggle PPT links update popup

    If `auto_update` is set to True, PPT is set to automatic update and the popup comes up whenever PPT is opened.
    If `auto_update` is set to False, PPT is set to manual update and the popup is no longer visible.
    """

    # Create a backup before proceeding
    if not overwrite:
        backup_path = _create_file_backup(ppt_path)

    # Create a temp directory to store our intermediate stuff
    with tempfile.TemporaryDirectory() as tmpdir:
        zfw = zipfile.ZipFile(f"{tmpdir}/temp.pptx", "w", compression=zipfile.ZIP_DEFLATED)
        with zipfile.ZipFile(ppt_path) as zfr:
            for ITEM in zfr.infolist():
                if ITEM.filename.startswith("ppt/"):
                    zfr.extract(ITEM.filename, tmpdir)
                else:
                    zfw.writestr(ITEM, zfr.open(ITEM).read(), compress_type=zipfile.ZIP_DEFLATED)

        # Search and replace string - CHARTS
        XML_FILES = glob.glob(f"{tmpdir}/ppt/charts/*.xml")
        search_str = f'<c:autoUpdate val="{int(not auto_update)}"/>'
        replace_with = f'<c:autoUpdate val="{int(auto_update)}"/>'
        _search_and_replace(search_str, replace_with, XML_FILES)

        # Search and replace string - TABLES
        XML_FILES = glob.glob(f"{tmpdir}/ppt/slides/*.xml")
        search_str = '<p:link' + (' updateAutomatic="1"' if not auto_update else '') + '/>'
        replace_with = '<p:link' + (' updateAutomatic="1"' if auto_update else '') + '/>'
        _search_and_replace(search_str, replace_with, XML_FILES)

        # Add the PPT folder we have updated to the archive
        for FILE_NAME in glob.glob(f"{tmpdir}\\ppt\\**", recursive=True):
            if os.path.isfile(FILE_NAME):
                FNP = FILE_NAME.split("\\")
                ZFNP = os.path.join(*FNP[FNP.index("ppt"):])
                zfw.write(FILE_NAME, ZFNP, compress_type=zipfile.ZIP_DEFLATED)

        # Close to save it
        zfw.close()

        # Replace the existing PPT
        shutil.move(f"{tmpdir}/temp.pptx", ppt_path)

def refresh_excel_external_connections(excel_path: str, debug: bool = False):
    """
    Uses xlwings to open an execl instance, refresh and close the excel file post update.
    """
    with xw.App(visible=debug, add_book=False) as app:
        # Open the workbook
        wb = app.books.open(excel_path)

        # Refresh all external data connections and wait until completed
        wb.api.RefreshAll()
        app.api.CalculateUntilAsyncQueriesDone()

        # Save the close the worksheet
        wb.save()
        wb.close()

def refresh_linked_plots_in_ppt(ppt_path: str, debug: bool = False):
    """
    Same as refreshing excel's external data connections, snippet refreshes the PPT plots linked via excel.
    Setting debug as True helps to see what goes wrong during the update.
    """

    # Open powerpoint
    PPTApp = win32com.client.Dispatch("PowerPoint.Application")
    PPTPresentation = PPTApp.Presentations.Open(ppt_path, WithWindow=debug)

    # Click on update links button
    PPTPresentation.UpdateLinks()

    # Save and close the PPT
    PPTPresentation.Save()
    PPTPresentation.Close()

    # Close the PPT Application
    PPTApp.Quit()

def extract_mqueries(excel_path: str, output_path: str) -> int:
    """
    Given an input excel, extracts the mqueries contained within, to the output
    path provided.
    """
    count = 0
    with zipfile.ZipFile(excel_path) as zfe:
        for excel_item in zfe.infolist():
            if re.match("customXml\/item\d+\.xml", excel_item.filename):
                tree = lxml.etree.fromstring(zfe.read(excel_item.filename))
                data_mashup_decoded = base64.b64decode(str(tree.text))
                count += _extract_excel_datamashup(data_mashup_decoded, output_path)

    return count

def update_mqueries(excel_path: str, mquery_paths: list[str], overwrite: bool = False) -> int:
    """
    Given an input excel, list of mquery files ('*.m' filepaths), the function
    updates the excel's datamashup contained within the path `customXml/item*.xml`
    """
    # Create a backup of the file before proceeding
    if not overwrite:
        backup_path = _create_file_backup(excel_path)

    # Convert Mquery FP to a dictionary to store our intermediate stuff
    mquery_fp: dict[str, str] = {os.path.basename(fp_): fp_ for fp_ in mquery_paths}

    # Create a temp directory to store our intermediate stuff
    touch_count = 0
    with tempfile.TemporaryDirectory() as tmpdir:
        zfw = zipfile.ZipFile(os.path.join(tmpdir, "temp.xlsx"), "w", compression=zipfile.ZIP_DEFLATED)
        with zipfile.ZipFile(excel_path) as zfr:
            for ITEM in zfr.infolist():
                if re.match("customXml\/item\d+\.xml", ITEM.filename):
                    custom_xml = zfr.open(ITEM).read()
                    tree = lxml.etree.fromstring(custom_xml)
                    data_mashup_decoded = base64.b64decode(str(tree.text))
                    data_mashup_updated = _update_excel_datamashup(data_mashup_decoded, mquery_fp)
                    tree.text = base64.b64encode(data_mashup_updated).decode("utf-8")
                    zfw.writestr(ITEM, lxml.etree.tostring(tree), compress_type=zipfile.ZIP_DEFLATED)
                    touch_count += 1
                else:
                    zfw.writestr(ITEM, zfr.open(ITEM).read(), compress_type=zipfile.ZIP_DEFLATED)

        # Close to save it
        zfw.close()

        # Replace the existing PPT
        shutil.move(f"{tmpdir}/temp.xlsx", excel_path)

    return touch_count

def extract_ppt(ppt_path: str, extract_path: str = "") -> int:
    """
    Extract PPT as a zip to path. If `extract_path` is not provided, creates a 'tmp'
    folder in the same directory as `ppt_path` and extracts there.
    """

    if not extract_path:
        extract_path = os.path.join(os.path.dirname(ppt_path), "tmp")

    files_extracted = 0
    with zipfile.ZipFile(ppt_path) as zfr:
        for ITEM in zfr.infolist():
            files_extracted += 1
            zfr.extract(ITEM.filename, extract_path)

    return files_extracted

def rezip_ppt(extract_path: str, ppt_path: str) -> None:
    """
    Given a path containing the unzipped PPT contents, zips and generates a PPT
    """
    with zipfile.ZipFile(ppt_path, "w", compression=zipfile.ZIP_DEFLATED) as zfw:
        for dirpath, dirnames, filenames in os.walk(extract_path):
            relative_dirpath = os.path.relpath(dirpath, extract_path)
            for filename in filenames:
                zfw.write(os.path.join(dirpath, filename), os.path.join(relative_dirpath, filename), compress_type=zipfile.ZIP_DEFLATED)

def update_ppt_plot_cache(extract_path: str) -> int:
    """
    Given a PPT extract path containing embedded excel files that are out of sync
    with the numCache, strCache - overwrites the cache directly from embeded
    excel charts.
    """
    update_count = 0
    chart_xml_path = os.path.join(extract_path, "ppt", "charts")
    rel_xml_path = os.path.join(extract_path, "ppt", "charts", "_rels")
    for filename in os.listdir(chart_xml_path):
        if filename.startswith("chart") and filename.endswith(".xml"):
            update_count += 1
            filepath = os.path.join(extract_path, "ppt", "charts", filename)
            rel_filepath = os.path.join(extract_path, "ppt", "charts", "_rels", f"{filename}.rels")

            # Get the embedded excel path from .rels file
            rel_filepath_xml = lxml.etree.parse(rel_filepath)
            relationship_element = rel_filepath_xml.find("*[@Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/package']")
            embed_excel_path = ""
            if relationship_element is not None:
                embed_excel_path = os.path.normpath(os.path.join(extract_path, "ppt", str(relationship_element.attrib['Target'])[3:]))
            assert embed_excel_path, f"{embed_excel_path} doesn't exist for {filepath}"

            # Read embed file chart xml and ppt xml side by side
            with zipfile.ZipFile(embed_excel_path, "r") as zfr:
                embed_chart_xml = lxml.etree.fromstring(zfr.read("xl/charts/chart1.xml"))
            ppt_chart_xml = lxml.etree.parse(filepath)

            # Replace the cache in current file from the embed excel cache
            for ppt_cache, embed_cache in zip(
                    ppt_chart_xml.xpath("//*[contains(local-name(), 'Cache')]"),
                    embed_chart_xml.xpath("//*[contains(local-name(), 'Cache')]")
                ):
                ppt_cache.getparent().replace(ppt_cache, embed_cache)

            # Overwrite the file changes
            ppt_chart_xml.write(filepath)

    return update_count

def update_embedded_excel(
        excel_instance: xw.App, embed_file_path: str, data_df: pd.DataFrame, *,
        sheet_name: str = 'data', fill_range: str = 'A1', file_sensitivity_id: str = '',
        fill_range_format: str = "[>=1]#,##0\%;[<1]0.0\%"
    ) -> None:
    """
    We take in an instance of excel and not create it inside a function so that we can have all
    the embed file updates done in one shot and not have to close and respawn a new excel instance
    for every file.

    Overwrite existing data using xlwings so that data is updated
    in real time (cannot do with pandas or openpyxl)
    """
    # Open the excel file
    wb = excel_instance.books.open(embed_file_path)

    # Add label before saving (optional)
    if file_sensitivity_id:
        labelinfo = wb.api.SensitibityLabel.CreateLabelInfo()
        labelinfo.AssignmentMethod = 2
        labelinfo.Justification = 'init'
        labelinfo.LabelId = file_sensitivity_id
        wb.api.SensitibityLabel.SetLabel(labelinfo, labelinfo)

    # Write the data to sheet 'data', starting from 1st cell
    wb.sheets[sheet_name].range(fill_range).options(index=False).value = data_df
    wb.save()
    wb.close()

def update_ppt_table(
        ppt_path: str, table_df: pd.DataFrame, *,
        slide_id: int, shape_id: int, start_coord: tuple[int, int],
        strides: tuple[int, int] = (1, 1), include_df_header: bool = True
    ) -> None:
    """
    Starts writing to the PPT table from start_coord provided.

    If include_df_header is set to true, first row is filled from df.columns, otherwise
    we start filling values from first row of table_df.

    To find the start_coord, use something like this:
    ```
    for i, row in enumerate(pptx.Presentation(TEMPLATE_PATH).slides[1].shapes[4].table.rows):
        for j, cell in enumerate(row.cells):
            print(i, j, cell.text)
    ```

    Stride is useful when the cells we are filling in spans across multiple rows / columns.
    Only useful when the stride is fixed for a specific range that we are trying to fill.
    For variable srides, it is better to manually use pptx to fill the table.
    """

    # Read the PPT file
    prs = pptx.Presentation(ppt_path)

    # Write to PPT Table
    table = prs.slides[slide_id].shapes[shape_id].table
    i_start, j_start = start_coord

    for i in range(table_df.shape[0] + (1 if include_df_header else 0)):
        for j in range(table_df.shape[1]):
            # Replace to make sure formatting is unaffected
            if include_df_header:
                replace = table_df.iloc[i - 1, j] if i > 0 else table_df.columns[j]
            else:
                replace = table_df.iloc[i, j]

            # Translate pandsa cell coord to PPT table cell coord
            ppt_cell_coord: tuple[int, int] = (i_start + (i * strides[0]), j_start + (j * strides[1]))
            assert not table.cell(*ppt_cell_coord).is_spanned, "tried to fill a hidden cell, please check `strides` and `start_coord` values provided."

            # Replace run text
            table.cell(*ppt_cell_coord).text_frame.paragraphs[0].runs[0].text = replace

    # Save PPT, overwrite existing
    prs.save(ppt_path)

def update_ppt_textboxes(ppt_path: str, textboxes: list[tuple[int, int, int, int, str]]) -> None:
    """
    Given a tuple of slide_ids, shape_ids, paragraph_ids, run_ids and text_content:
        Updates the text of these textboxes

    Some trial and error is required to figure out the correct_ids pointing to the text content
    we desire to update.
    """
    prs = pptx.Presentation(ppt_path)
    for slide_id, shape_id, para_id, run_id, text in textboxes:
        prs.slides[slide_id].shapes[shape_id].text_frame.paragraphs[para_id].runs[run_id].text = text
    prs.save(ppt_path)
